
# # from pytube import extract
# # from youtube_transcript_api import YouTubeTranscriptApi
# # from youtube_transcript_api.formatters import TextFormatter, JSONFormatter

from langchain_core.documents import Document
from langchain_community.document_loaders.base import BaseBlobParser
from langchain_community.document_loaders.blob_loaders import Blob

import pandas as pd
import json
from typing import Iterator
from abc import ABC
import pptx
from io import BytesIO

# # class YoutubeParser(BaseBlobParser, ABC):
# #     def __init__(self):
# #         self.formatter = TextFormatter()

# #     def lazy_parse(self, blob: Blob) -> Iterator[Document]:
# #         video_id = extract.video_id(blob.source)

# #         transcript = YouTubeTranscriptApi.get_transcripts([video_id], languages=["en", "it"], preserve_formatting=True)
# #         text = self.formatter.format_transcript(transcript[0][video_id])

# #         yield Document(page_content=text, metadata={})

class TableParser(BaseBlobParser, ABC):

    def lazy_parse(self, blob: Blob) -> Iterator[Document]:
        with blob.as_bytes_io() as file:
            if blob.mimetype == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":

                # Read the Excel file - df is a dictionary of DataFrames (one per sheet)
                df_dict = pd.read_excel(file, sheet_name=None)
                all_records = []
                # Process each sheet
                for sheet_name, sheet_df in df_dict.items():
                    # Convert each sheet's DataFrame to records and add to our list
                    sheet_records = sheet_df.to_dict("records")
                    all_records.extend(sheet_records)

            elif blob.mimetype == "text/csv":
                # For CSV files, use read_csv instead of read_excel
                df = pd.read_csv(file)
                all_records = df.to_dict("records")
            
        yield Document(page_content=json.dumps(all_records), metadata={})

class PowerPointParser(BaseBlobParser, ABC):
    
    def lazy_parse(self, blob: Blob) -> Iterator[Document]:
        # Accept multiple PowerPoint MIME types
        pptx_mime_types = [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
            "application/vnd.ms-powerpoint",  # .ppt
            "application/powerpoint"  # Alternative for .ppt
        ]
        
        if blob.mimetype not in pptx_mime_types:
            raise ValueError(f"Unsupported mime type: {blob.mimetype}")
        
        with blob.as_bytes_io() as file_obj:
            presentation = pptx.Presentation(file_obj)
            
            # Extract text from all slides
            all_text = []
            slide_contents = {}
            
            for i, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Get slide title if available
                title = ""
                for shape in slide.shapes:
                    # Check if shape has text attribute and if it's a title shape
                    if hasattr(shape, "text") and hasattr(shape, "is_title") and shape.is_title:
                        title = shape.text
                        break
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                # Join all text from this slide
                slide_content = "\n".join(slide_text)
                all_text.append(slide_content)
                
                # Add to slide_contents dictionary
                slide_contents[f"Slide {i}"] = {
                    "title": title,
                    "content": slide_content
                }
            
            # Join all text from all slides
            full_text = "\n\n".join(all_text)
            
            yield Document(page_content=full_text, metadata={"slide_contents": slide_contents})

class EmailParser(BaseBlobParser, ABC):
    
    def lazy_parse(self, blob: Blob) -> Iterator[Document]:
        # Accept email MIME types
        email_mime_types = [
            "message/rfc822",  # .eml
            "application/vnd.ms-outlook",  # .msg
            "application/octet-stream"  # Sometimes used for email files
        ]
        
        if blob.mimetype not in email_mime_types and not (blob.path.lower().endswith('.eml') or blob.path.lower().endswith('.msg')):
            raise ValueError(f"Unsupported mime type: {blob.mimetype}")
        
        with blob.as_bytes_io() as file_obj:
            # For .eml files (RFC822 format)
            if blob.mimetype == "message/rfc822" or blob.path.lower().endswith('.eml'):
                import email
                from email import policy
                
                try:
                    from bs4 import BeautifulSoup
                    has_bs4 = True
                except ImportError:
                    has_bs4 = False
                
                # Parse the email
                msg = email.message_from_binary_file(file_obj, policy=policy.default)
                
                # Extract only essential headers
                subject = msg.get("Subject", "")
                sender = msg.get("From", "")
                recipients = msg.get("To", "")
                cc = msg.get("Cc", "")
                
                # Extract body content
                body = ""
                
                # Handle multipart messages
                if msg.is_multipart():
                    for part in msg.iter_parts():
                        content_type = part.get_content_type()
                        if content_type == "text/plain":
                            body += part.get_content() + "\n\n"
                        elif content_type == "text/html" and not body and has_bs4:
                            # Extract text from HTML content
                            html_content = part.get_content()
                            soup = BeautifulSoup(html_content, 'html.parser')
                            body += soup.get_text(separator='\n') + "\n\n"
                else:
                    # Handle single part messages
                    content_type = msg.get_content_type()
                    if content_type == "text/plain":
                        body = msg.get_content()
                    elif content_type == "text/html" and has_bs4:
                        html_content = msg.get_content()
                        soup = BeautifulSoup(html_content, 'html.parser')
                        body = soup.get_text(separator='\n')
            
            # For .msg files (Outlook format)
            elif blob.mimetype == "application/vnd.ms-outlook" or blob.path.lower().endswith('.msg'):
                import extract_msg
                
                # Reset file pointer to beginning
                file_obj.seek(0)
                
                # Save to a temporary file since extract_msg needs a file path
                import tempfile
                import os
                
                with tempfile.NamedTemporaryFile(delete=False) as temp:
                    temp.write(file_obj.read())
                    temp_path = temp.name
                
                try:
                    # Parse the .msg file
                    outlook_msg = extract_msg.Message(temp_path)
                    
                    # Extract only essential information
                    subject = outlook_msg.subject
                    sender = outlook_msg.sender
                    recipients = outlook_msg.to
                    cc = outlook_msg.cc
                    
                    # Extract body
                    body = outlook_msg.body
                    
                    # Close the message
                    outlook_msg.close()
                finally:
                    # Clean up the temporary file
                    os.unlink(temp_path)
            else:
                raise ValueError(f"Unsupported email format: {blob.mimetype}")
            
            # Format the essential content in a clean, readable format
            email_parts = []
            if sender:
                email_parts.append(f"Da: {sender}")
            if recipients:
                email_parts.append(f"A: {recipients}")
            if cc:
                email_parts.append(f"CC: {cc}")
            if subject:
                email_parts.append(f"Oggetto: {subject}")
            if body:
                email_parts.append("\n" + body)
            
            # Join all parts with newlines
            full_content = "\n".join(email_parts)
            
            yield Document(page_content=full_content, metadata={})

# class JSONParser(BaseBlobParser, ABC):

#     def lazy_parse(self, blob: Blob) -> Iterator[Document]:

#         with blob.as_bytes_io() as file:
#             text = json.load(file)

#         yield Document(page_content=text, metadata={})

